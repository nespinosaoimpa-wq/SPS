from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r"c:\Users\Grupo 5\Desktop\SPS Prototipo\power de presentación básico SISTEMA DE APOYO Y CONTROL AL SERVICIO PROFESIONAL.pptx"
prs = Presentation(pptx_path)

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    print(f"Title: {slide.shapes.title.text if slide.shapes.title else 'No Title'}")
    for s_idx, shape in enumerate(slide.shapes):
        print(f"Shape {s_idx+1}: Name='{shape.name}', Type={shape.shape_type}")
        if shape.is_placeholder:
            ph = shape.placeholder_format
            print(f"  Is Placeholder: idx={ph.idx}, type={ph.type}")
        if shape.has_text_frame:
            print(f"  Has Text Frame. Text: '{shape.text_frame.text.strip()}'")
        if shape.has_table:
            print(f"  Has Table.")
        # Check if shape contains image
        try:
            if hasattr(shape, 'image'):
                print(f"  Has image attribute! Ext={shape.image.ext}")
        except Exception as e:
            print(f"  Image check error: {e}")
        # Check if shape is a picture or graphic frame
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"  It is a PICTURE shape.")
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            print(f"  It is a PLACEHOLDER shape.")
