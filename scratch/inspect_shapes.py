from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r"c:\Users\Grupo 5\Desktop\SPS Prototipo\power de presentación básico SISTEMA DE APOYO Y CONTROL AL SERVICIO PROFESIONAL.pptx"
prs = Presentation(pptx_path)

def inspect_shape(shape, depth=0):
    indent = "  " * depth
    print(f"{indent}- Shape: {shape.name}, Type: {shape.shape_type}")
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            inspect_shape(s, depth + 1)
    else:
        if shape.has_text_frame:
            print(f"{indent}  Text: {shape.text_frame.text.strip()[:100]}")
        if shape.has_table:
            print(f"{indent}  Table rows: {len(shape.table.rows)}")

slides_list = list(prs.slides)
for i, slide in enumerate(slides_list[:3]):
    print(f"\nSlide {i+1}:")
    for shape in slide.shapes:
        inspect_shape(shape)
