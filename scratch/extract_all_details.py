import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r"c:\Users\Grupo 5\Desktop\SPS Prototipo\power de presentación básico SISTEMA DE APOYO Y CONTROL AL SERVICIO PROFESIONAL.pptx"
prs = Presentation(pptx_path)

out_dir = r"c:\Users\Grupo 5\Desktop\SPS Prototipo\scratch\extracted_images"
os.makedirs(out_dir, exist_ok=True)

print(f"Total slides: {len(prs.slides)}")

def extract_from_shape(shape, slide_num, shape_num_ref):
    details = []
    
    # Text extraction
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
    
    # Table extraction
    table_data = []
    if shape.has_table:
        for r_idx, row in enumerate(shape.table.rows):
            row_cells = [cell.text.strip() for cell in row.cells]
            table_data.append(" | ".join(row_cells))
            
    # Image extraction - check if has attribute image
    image_saved = None
    if hasattr(shape, 'image') and shape.image is not None:
        try:
            image = shape.image
            ext = image.ext
            image_bytes = image.blob
            filename = f"slide_{slide_num}_shape_{shape_num_ref[0]}.{ext}"
            filepath = os.path.join(out_dir, filename)
            with open(filepath, "wb") as img_file:
                img_file.write(image_bytes)
            image_saved = filepath
            shape_num_ref[0] += 1
        except Exception as e:
            print(f"Error extracting image from shape {shape.name} in slide {slide_num}: {e}")
        
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            child_details = extract_from_shape(child, slide_num, shape_num_ref)
            details.extend(child_details)
        return details

    return [{
        "name": shape.name,
        "type": str(shape.shape_type),
        "text": text,
        "table": table_data,
        "image": image_saved
    }]

slides_info = []
for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    shape_num_ref = [1]
    slide_data = {
        "num": slide_num,
        "title": slide.shapes.title.text if slide.shapes.title else f"Slide {slide_num}",
        "shapes": []
    }
    for shape in slide.shapes:
        shape_data = extract_from_shape(shape, slide_num, shape_num_ref)
        if shape_data:
            slide_data["shapes"].extend(shape_data)
    slides_info.append(slide_data)

# Let's write the details to a markdown report
report_path = r"c:\Users\Grupo 5\Desktop\SPS Prototipo\scratch\slide_analysis.md"
with open(report_path, "w", encoding="utf-8") as f:
    f.write("# Detalle Analítico de Diapositivas\n\n")
    for s in slides_info:
        f.write(f"## Diapositiva {s['num']}: {s['title']}\n\n")
        
        has_content = False
        for shape in s["shapes"]:
            if shape["text"] and shape["text"] != s["title"]:
                f.write(f"### {shape['name']} ({shape['type']})\n")
                f.write(f"**Texto:** {shape['text']}\n\n")
                has_content = True
            if shape["table"]:
                f.write(f"### {shape['name']} (Tabla)\n")
                for row in shape["table"]:
                    f.write(f"- {row}\n")
                f.write("\n")
                has_content = True
            if shape["image"]:
                f.write(f"### {shape['name']} (Imagen)\n")
                # Using standard markdown format for absolute path links
                f.write(f"Guardada como: [{os.path.basename(shape['image'])}](file:///{shape['image'].replace('\\', '/')})\n\n")
                has_content = True
                
        if not has_content:
            f.write("*Esta diapositiva solo contiene el título u otros elementos no textuales.*\n\n")
        f.write("---\n\n")

print(f"Analysis complete. Report written to {report_path}")
