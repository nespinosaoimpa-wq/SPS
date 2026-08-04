import os
from pptx import Presentation

pptx_path = r"c:\Users\Grupo 5\Desktop\SPS Prototipo\power de presentación básico SISTEMA DE APOYO Y CONTROL AL SERVICIO PROFESIONAL.pptx"
prs = Presentation(pptx_path)

out_path = r"c:\Users\Grupo 5\Desktop\SPS Prototipo\scratch\pptx_content.md"

with open(out_path, "w", encoding="utf-8") as f:
    f.write("# Contenido de la Presentación PowerPoint\n\n")
    for i, slide in enumerate(prs.slides):
        f.write(f"## Diapositiva {i+1}\n\n")
        
        # Try to find a title
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text
            f.write(f"### Título: {title}\n\n")
        
        # Extract text from all shapes
        text_runs = []
        for shape in slide.shapes:
            # Check for standard text frame
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    p_text = paragraph.text.strip()
                    if p_text and p_text != title:
                        text_runs.append(p_text)
            # Check for table
            elif shape.has_table:
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    text_runs.append(" | ".join(row_text))
                    
        if text_runs:
            f.write("**Contenido:**\n")
            for run in text_runs:
                f.write(f"- {run}\n")
            f.write("\n")
        else:
            f.write("*Sin contenido de texto*\n\n")
            
        f.write("---\n\n")

print("Extract complete! Written to scratch/pptx_content.md")
